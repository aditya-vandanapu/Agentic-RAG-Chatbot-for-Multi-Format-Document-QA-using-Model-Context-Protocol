import docx2txt
import csv
from PyPDF2 import PdfReader
from pptx import Presentation

def parse_txt(file): return file.read().decode("utf-8")

def parse_docx(file): return docx2txt.process(file)

def parse_pdf(file):
    reader = PdfReader(file)
    texts = []
    for i, page in enumerate(reader.pages):
        content = page.extract_text()
        if content:
            texts.append((f"page_{i+1}", content.strip()))
    return texts

def parse_csv(file):
    content = file.read().decode("utf-8").splitlines()
    rows = list(csv.reader(content))
    return [("csv", "\n".join([", ".join(row) for row in rows]))]

def parse_pptx(file):
    prs = Presentation(file)
    texts = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        texts.append((f"slide_{i+1}", "\n".join(slide_text)))
    return texts

def parse_file(file, filename):
    ext = filename.split(".")[-1].lower()
    if ext == "txt" or ext == "md": return [("text", parse_txt(file))]
    if ext == "pdf": return parse_pdf(file)
    if ext == "docx": return [("docx", parse_docx(file))]
    if ext == "csv": return parse_csv(file)
    if ext == "pptx": return parse_pptx(file)
    return []
